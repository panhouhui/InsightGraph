import logging
import os
from pathlib import Path
import tempfile
import zipfile
import chardet
import pyzipper
from langchain_core.documents import Document
from langchain_core.document_loaders import BaseLoader
from src.shared.llm_graph_builder_exception import LLMGraphBuilderException

SUPPORTED_ARCHIVE_EXTENSIONS = {
    ".pdf",
    ".jpeg",
    ".jpg",
    ".png",
    ".svg",
    ".html",
    ".docx",
    ".txt",
    ".pptx",
    ".xls",
    ".md",
    ".xlsx",
}


def get_int_from_env(key_name, default_value):
    value = os.getenv(key_name)
    if value in (None, ""):
        return default_value
    return int(value)

class ListLoader(BaseLoader):
    """
    A wrapper to make a list of Documents compatible with BaseLoader.
    """
    def __init__(self, documents):
        self.documents = documents

    def load(self):
        """
        Returns the list of documents.
        """
        return self.documents

def detect_encoding(file_path):
    """
    Detects the file encoding to avoid UnicodeDecodeError.

    Args:
        file_path (str or Path): Path to the file.

    Returns:
        str: Detected encoding (default "utf-8" if not found).
    """
    with open(file_path, 'rb') as f:
        raw_data = f.read(4096)
        result = chardet.detect(raw_data)
        return result['encoding'] or "utf-8"

def _is_zip_member_encrypted(info):
    """Return True for traditional encrypted or AES encrypted zip members."""
    is_traditional_encrypted = bool(info.flag_bits & 0x1)
    is_aes_encrypted = b"\x01\x99" in info.extra or b"\x99\x01" in info.extra
    return is_traditional_encrypted or is_aes_encrypted


def _safe_zip_target(extract_dir, member_name):
    target_path = (Path(extract_dir) / member_name).resolve()
    extract_root = Path(extract_dir).resolve()
    try:
        target_path.relative_to(extract_root)
    except ValueError as exc:
        raise LLMGraphBuilderException(f"Unsafe archive entry path: {member_name}") from exc
    return target_path


def _load_pages_from_extracted_file(file_path, display_name):
    loader, encoding_flag = load_document_content(file_path)
    file_extension = Path(file_path).suffix.lower()
    if file_extension == ".pdf" or (file_extension == ".txt" and encoding_flag):
        pages = loader.load()
    else:
        pages = get_pages_with_page_numbers(loader.load())

    for page in pages:
        page.metadata = {
            **page.metadata,
            "source": display_name,
            "filename": display_name,
        }
    return pages


def _document_from_text(text, source, filetype):
    return Document(
        page_content=text,
        metadata={
            "source": source,
            "filename": Path(source).name,
            "filetype": filetype,
            "page_number": 1,
        },
    )


def _load_docx(file_path):
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    parts = []
    parts.extend(paragraph.text for paragraph in doc.paragraphs if paragraph.text and paragraph.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return [_document_from_text("\n".join(parts), str(file_path), "docx")]


def _load_xlsx(file_path):
    from openpyxl import load_workbook

    workbook = load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None and str(value).strip()]
                if values:
                    parts.append(" | ".join(values))
    finally:
        workbook.close()
    return [_document_from_text("\n".join(parts), str(file_path), "xlsx")]


def _load_pptx(file_path):
    from pptx import Presentation

    presentation = Presentation(file_path)
    parts = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"Slide {slide_index}\n" + "\n".join(slide_parts))
    return [_document_from_text("\n".join(parts), str(file_path), "pptx")]


def _load_office_document(file_path):
    file_extension = Path(file_path).suffix.lower()
    if file_extension == ".docx":
        return _load_docx(file_path)
    if file_extension == ".xlsx":
        return _load_xlsx(file_path)
    if file_extension == ".pptx":
        return _load_pptx(file_path)
    return None


def load_documents_from_zip(file_path, archive_password=None):
    """
    Load supported document files from a zip archive.

    Encrypted archives require archive_password. AES encrypted archives are read with
    pyzipper because Python's stdlib zipfile cannot decrypt them.
    """
    max_files = get_int_from_env("ARCHIVE_MAX_FILES", 100)
    max_uncompressed_bytes = get_int_from_env("ARCHIVE_MAX_UNCOMPRESSED_BYTES", 200 * 1024 * 1024)
    password_bytes = archive_password.encode("utf-8") if archive_password else None

    try:
        with pyzipper.AESZipFile(file_path) as archive:
            infos = [
                info for info in archive.infolist()
                if not info.is_dir() and not Path(info.filename).name.startswith(".")
            ]
            supported_infos = [
                info for info in infos
                if Path(info.filename).suffix.lower() in SUPPORTED_ARCHIVE_EXTENSIONS
            ]

            if not supported_infos:
                raise LLMGraphBuilderException("Archive does not contain any supported document files")
            if len(supported_infos) > max_files:
                raise LLMGraphBuilderException(f"Archive contains too many supported files. Limit is {max_files}")

            total_size = sum(info.file_size for info in supported_infos)
            if total_size > max_uncompressed_bytes:
                raise LLMGraphBuilderException(
                    f"Archive uncompressed size exceeds the limit of {max_uncompressed_bytes} bytes"
                )

            if any(_is_zip_member_encrypted(info) for info in supported_infos):
                if not password_bytes:
                    raise LLMGraphBuilderException("Archive is encrypted. Please provide the archive password")
                archive.setpassword(password_bytes)

            pages = []
            with tempfile.TemporaryDirectory(prefix="llm_graph_builder_zip_") as extract_dir:
                for info in supported_infos:
                    target_path = _safe_zip_target(extract_dir, info.filename)
                    target_path.parent.mkdir(parents=True, exist_ok=True)
                    try:
                        with archive.open(info, pwd=password_bytes) as source_file, open(target_path, "wb") as target_file:
                            target_file.write(source_file.read())
                    except (RuntimeError, pyzipper.zipfile.BadZipFile) as exc:
                        raise LLMGraphBuilderException("Invalid archive password or encrypted archive cannot be read") from exc

                    display_name = f"{Path(file_path).name}::{info.filename}"
                    pages.extend(_load_pages_from_extracted_file(target_path, display_name))
    except zipfile.BadZipFile as exc:
        raise LLMGraphBuilderException("Archive is not a valid zip file") from exc

    return pages


def load_document_content(file_path):
    """
    Loads document content from a file, handling PDFs and text encoding.

    Args:
        file_path (str or Path): Path to the file.

    Returns:
        tuple: (loader, encoding_flag)
            loader: Document loader instance.
            encoding_flag (bool): True if non-UTF-8 encoding was used for .txt files.
    """
    file_extension = Path(file_path).suffix.lower()
    encoding_flag = False
    if file_extension == '.pdf':
        from langchain_community.document_loaders import PyMuPDFLoader

        loader = PyMuPDFLoader(file_path)
        return loader, encoding_flag
    office_pages = _load_office_document(file_path)
    if office_pages is not None:
        return ListLoader(office_pages), True
    if file_extension == ".txt":
        encoding = detect_encoding(file_path)
        logging.info("Detected encoding for text file: %s", encoding)
        if encoding.lower() == "utf-8":
            from langchain_community.document_loaders import UnstructuredFileLoader

            loader = UnstructuredFileLoader(file_path, mode="elements", autodetect_encoding=True)
            return loader, encoding_flag
        with open(file_path, encoding=encoding, errors="replace") as f:
            content = f.read()
        loader = ListLoader([Document(page_content=content, metadata={"source": file_path})])
        encoding_flag = True
        return loader, encoding_flag
    from langchain_community.document_loaders import UnstructuredFileLoader

    loader = UnstructuredFileLoader(file_path, mode="elements", autodetect_encoding=True)
    return loader, encoding_flag

def get_documents_from_file_by_path(file_path, file_name, archive_password=None):
    """
    Loads documents from a file by its path and returns file name, pages, and extension.

    Args:
        file_path (str or Path): Path to the file.
        file_name (str): Name of the file.

    Returns:
        tuple: (file_name, pages, file_extension)

    Raises:
        Exception: If file does not exist or reading fails.
    """
    file_path = Path(file_path)
    if not file_path.exists():
        logging.info('File %s does not exist', file_name)
        raise Exception(f'File {file_name} does not exist')
    logging.info('file %s processing', file_name)
    try:
        file_extension = file_path.suffix.lower()
        if file_extension == ".zip":
            pages = load_documents_from_zip(file_path, archive_password)
        else:
            loader, encoding_flag = load_document_content(file_path)
        if file_extension == ".pdf" or (file_extension == ".txt" and encoding_flag):
            pages = loader.load()
        elif file_extension != ".zip":
            unstructured_pages = loader.load()
            pages = get_pages_with_page_numbers(unstructured_pages)
    except LLMGraphBuilderException:
        raise
    except Exception as exc:
        raise Exception(f'Error while reading the file content or metadata, {exc}')
    return file_name, pages, file_extension

def get_pages_with_page_numbers(unstructured_pages):
    """
    Groups unstructured pages into logical pages with page numbers and metadata.

    Args:
        unstructured_pages (list): List of Document objects.

    Returns:
        list: List of Document objects with page numbers and metadata.
    """
    pages = []
    page_number = 1
    page_content = ''
    metadata = {}
    for idx, page in enumerate(unstructured_pages):
        if 'page_number' in page.metadata:
            if page.metadata['page_number'] == page_number:
                page_content += page.page_content
                metadata = {
                    'source': page.metadata['source'],
                    'page_number': page_number,
                    'filename': page.metadata['filename'],
                    'filetype': page.metadata['filetype']
                }
            if page.metadata['page_number'] > page_number:
                page_number += 1
                pages.append(Document(page_content=page_content))
                page_content = ''
            if page == unstructured_pages[-1]:
                pages.append(Document(page_content=page_content))
        elif page.metadata.get('category') == 'PageBreak' and page != unstructured_pages[0]:
            page_number += 1
            pages.append(Document(page_content=page_content, metadata=metadata))
            page_content = ''
            metadata = {}
        else:
            page_content += page.page_content
            metadata_with_custom_page_number = {
                'source': page.metadata['source'],
                'page_number': 1,
                'filename': page.metadata['filename'],
                'filetype': page.metadata['filetype']
            }
            if page == unstructured_pages[-1]:
                pages.append(Document(page_content=page_content, metadata=metadata_with_custom_page_number))
    return pages
